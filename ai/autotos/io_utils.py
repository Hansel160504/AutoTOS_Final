"""Text I/O: file extraction, chunking, and small string helpers.

Pure functions where possible; chunking has a bounded LRU cache.

v21 update: added _strip_filler_phrases() to remove academic scaffolding
("in the context of X", "within the Y structure", "based on the lesson",
etc.) from generated question stems. Called from strip_question_prefix()
so it runs on every question via normalize_question().
"""
from __future__ import annotations

import base64
import hashlib
import logging
import os
import re
from collections import OrderedDict
from io import BytesIO
from typing import Optional

from .config import CHUNK_CACHE_MAX, CHUNK_OVERLAP, CHUNK_SIZE, MAX_RETURN, CACHE_DIR

logger = logging.getLogger(__name__)

# Lazy-imported in helpers so tests don't pull heavy deps unnecessarily.

# ── String helpers ────────────────────────────────────────────────────
_WS_RE = re.compile(r"\s+")
_ANSWER_PREFIX_RE     = re.compile(r"^(answer\s*[:\-]\s*)", re.IGNORECASE)
_TF_QUESTION_PREFIX   = re.compile(r"^(true\s+or\s+false\s*[:\-]\s*)", re.IGNORECASE)
_LEADING_DIGIT_RE     = re.compile(r"\b\d+\s+(?=[A-Z])")
_FILL_BLANK_RE        = re.compile(r"_{3,}")
_ABBR_RE = re.compile(
    r"\b(e\.g|i\.e|etc|vs|Mr|Mrs|Dr|Prof|Sr|Jr|St|approx|fig|no)\.\s",
    re.IGNORECASE,
)

# ── Filler-phrase stripping ───────────────────────────────────────────
# These patterns wrap questions in academic-sounding scaffolding without
# adding information. Examples we strip:
#   "What does X refer to in the context of Y?" → "What does X refer to?"
#   "In the context of X, what is Y?"           → "What is Y?"
#   "What is the role of X within the Y system?" → "What is the role of X?"
#   "What is X based on the lesson?"             → "What is X?"
#
# Order matters — patterns are tried sequentially, and earlier ones may
# strip text that later ones would have caught. Conservative: only fires
# on text >20 chars, leaves substantive uses alone.

# Mid/end-of-sentence filler — preserves trailing punctuation
_FILLER_PHRASE_RES = [
    # "... in the context of <topic>?"
    re.compile(
        r",?\s+in the context of\s+[^?.!]{3,80}([?.!])",
        re.IGNORECASE,
    ),
    # "... within the <topic> structure/system/format/etc?"
    re.compile(
        r",?\s+within the\s+[^?.!]{3,60}\s+(?:structure|system|format|architecture|context|framework|environment)([?.!])",
        re.IGNORECASE,
    ),
    # "... based on (the/this/that) <noun>?"
    re.compile(
        r",?\s+based on\s+(?:the|this|that)\s+(?:lesson|context|material|content|reading|text|document|source)([?.!])",
        re.IGNORECASE,
    ),
    # "... according to (the/this/that) <noun>?"
    re.compile(
        r",?\s+according to\s+(?:the|this|that)\s+(?:lesson|context|material|content|reading|text|document|source|presentation)([?.!])",
        re.IGNORECASE,
    ),
    # "... as described/discussed/mentioned in <noun>?"
    re.compile(
        r",?\s+as (?:described|discussed|mentioned|outlined|stated|noted)\s+(?:in|by)\s+[^?.!]{3,40}([?.!])",
        re.IGNORECASE,
    ),
]

# Sentence-start filler — handled separately (replaced with empty string)
_FILLER_SENTENCE_START = re.compile(
    r"^in the context of\s+[^,]{3,80},\s+",
    re.IGNORECASE,
)


def _strip_filler_phrases(text: str) -> str:
    """Remove filler scaffolding from question stems.

    Conservative: only fires on text >20 chars. Preserves trailing
    punctuation. Re-capitalizes the first letter after stripping
    leading filler.
    """
    if not text or len(text) < 20:
        return text

    out = text

    # Strip sentence-start filler first
    out = _FILLER_SENTENCE_START.sub("", out)

    # Then strip mid/end filler — preserves the captured punctuation
    for pat in _FILLER_PHRASE_RES:
        out = pat.sub(lambda m: m.group(1) if m.lastindex else "?", out)

    # Clean up artifacts from substitution
    out = re.sub(r"\s{2,}", " ", out)               # collapse double spaces
    out = re.sub(r"\s+([?.!,])", r"\1", out)        # remove space before punctuation
    out = re.sub(r"^[,\s]+", "", out)               # leading commas/spaces
    if out and out[0].islower():                    # re-capitalize first letter
        out = out[0].upper() + out[1:]

    return out.strip()


def clean_text(txt) -> str:
    """Collapse whitespace and trim. Coerces non-strings safely."""
    if txt is None:
        return ""
    if not isinstance(txt, str):
        try:
            txt = str(txt)
        except Exception:
            return ""
    return _WS_RE.sub(" ", txt).strip()


def is_fill_in_blank(question: str) -> bool:
    return bool(question) and bool(_FILL_BLANK_RE.search(question))


def strip_question_prefix(text: str, *, is_open_ended: bool = False) -> str:
    """Strip auto-generated TF/numbered prefixes; remove filler scaffolding;
    uppercase first letter."""
    text = (text or "").strip()
    if not text:
        return text
    text = _LEADING_DIGIT_RE.sub("", text).strip()
    if not is_open_ended:
        text = _TF_QUESTION_PREFIX.sub("", text).strip()
        text = re.sub(r"^it is true that\s+", "", text, flags=re.IGNORECASE).strip()
    # Strip filler scaffolding ('in the context of X', 'within Y', etc.)
    text = _strip_filler_phrases(text)
    return text[:1].upper() + text[1:] if text else text


def truncate_answer_text(text: str, max_chars: int) -> str:
    """Single-sentence truncation for answer_text fields."""
    if not text:
        return ""
    t = _ANSWER_PREFIX_RE.sub("", _WS_RE.sub(" ", text).strip()).strip()
    sentences = re.split(r"(?<=[.!?])\s+", t)
    first = sentences[0].strip() if sentences else t
    if first and first[-1] not in ".!?":
        first += "."
    if len(first) > max_chars:
        cut = first[:max_chars]
        sp  = cut.rfind(" ")
        if sp > 0:
            cut = cut[:sp]
        return cut.rstrip(".,;:") + "..."
    return first


def truncate_open_answer(text: str, max_sentences: int = 4) -> str:
    """Return up to N complete sentences (handles abbreviations)."""
    if not text:
        return ""
    t = _ANSWER_PREFIX_RE.sub("", _WS_RE.sub(" ", text).strip()).strip()
    if not t:
        return ""
    masked = _ABBR_RE.sub(lambda m: m.group(0).replace(".", "\x00"), t)
    parts  = re.split(r"(?<=[.!?])\s+(?=[A-Z])|(?<=[.!?])$", masked)
    parts  = [p.replace("\x00", ".").strip() for p in parts if p.strip()]
    result = " ".join(parts[:max_sentences])
    if result and result[-1] not in ".!?":
        result += "."
    return result


# ── Chunking with bounded LRU ─────────────────────────────────────────
class _BoundedDict(OrderedDict):
    def __init__(self, maxlen: int) -> None:
        super().__init__()
        self._maxlen = maxlen

    def __setitem__(self, key, value):
        super().__setitem__(key, value)
        self.move_to_end(key)
        while len(self) > self._maxlen:
            self.popitem(last=False)


_chunk_cache: _BoundedDict = _BoundedDict(CHUNK_CACHE_MAX)


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Slide-window chunker preferring sentence then word boundaries."""
    if not text:
        return []
    text = text.strip()
    out: list[str] = []
    start = 0
    step = max(1, size - overlap)

    while start < len(text):
        end = min(start + size, len(text))
        if end < len(text):
            mid = start + size // 2
            sent_end = max(
                text.rfind(". ", start, end),
                text.rfind("! ", start, end),
                text.rfind("? ", start, end),
            )
            if sent_end > mid:
                end = sent_end + 1
            else:
                word_end = text.rfind(" ", start, end)
                if word_end > start:
                    end = word_end
        chunk = text[start:end].strip()
        if chunk:
            out.append(chunk)
        start += step
    return out


def get_chunks_for_text(full_text: str) -> list[str]:
    """Memoised chunking — bounded so repeated docs don't blow memory."""
    key = hashlib.md5(full_text[:4096].encode("utf-8", errors="ignore")).hexdigest()
    cached = _chunk_cache.get(key)
    if cached is not None:
        # Refresh recency
        _chunk_cache[key] = cached
        return cached
    chunks = chunk_text(full_text)
    _chunk_cache[key] = chunks
    logger.info("Chunked document: %d chars → %d chunks", len(full_text), len(chunks))
    return chunks


def find_best_chunk_idx(chunks: list[str], topic: str) -> int:
    """Pick the chunk most likely to discuss the topic."""
    if not chunks or not topic:
        return 0
    topic_lc = topic.lower()
    for i, ch in enumerate(chunks):
        if topic_lc in ch.lower():
            return i
    for word in topic_lc.split():
        if len(word) > 3:
            for i, ch in enumerate(chunks):
                if word in ch.lower():
                    return i
    return 0


# ── File extraction (PDF / DOCX / PPTX) ───────────────────────────────
def _extract_pdf(data: bytes) -> str:
    import fitz  # heavy import: lazy
    parts: list[str] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            try:
                txt = page.get_text("text")
                if txt and txt.strip():
                    parts.append(txt)
            except Exception:
                continue
    return " ".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    d = Document(BytesIO(data))
    return " ".join(p.text for p in d.paragraphs if p.text and p.text.strip())


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text)
    return " ".join(parts)


_EXTRACTORS = {"pdf": _extract_pdf, "docx": _extract_docx, "pptx": _extract_pptx}


def extract_text(file_bytes: bytes, filetype: str) -> str:
    extractor = _EXTRACTORS.get(filetype)
    try:
        if extractor:
            return clean_text(extractor(file_bytes))
        return clean_text(file_bytes.decode("utf-8", errors="ignore"))
    except Exception as exc:
        logger.warning("Extraction error (%s): %s", filetype, exc)
        return ""


def _filetype_from_ext(path: str) -> str:
    ext = (os.path.splitext(path)[1] or "").lower().lstrip(".")
    return {"pdf": "pdf", "docx": "docx", "doc": "docx",
            "pptx": "pptx", "ppt": "pptx"}.get(ext, "")


def _filetype_from_data_url_header(header: str) -> str:
    h = header.lower()
    if "pdf" in h:
        return "pdf"
    if "docx" in h or "wordprocessingml" in h or "msword" in h:
        return "docx"
    if "pptx" in h or "presentationml" in h or "ms-powerpoint" in h:
        return "pptx"
    return ""


def _sha256(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()


def extract_from_path(path: str, max_chars: int = MAX_RETURN) -> str:
    """Read a local file, extract text, cache by content hash."""
    if not path or not os.path.exists(path):
        return ""
    try:
        with open(path, "rb") as f:
            data = f.read()
    except OSError as exc:
        logger.warning("Read failed for %s: %s", path, exc)
        return ""

    digest = _sha256(data)
    cache_file = os.path.join(CACHE_DIR, f"{digest}.txt")
    if os.path.exists(cache_file):
        try:
            with open(cache_file, "r", encoding="utf-8") as f:
                return clean_text(f.read())[:max_chars]
        except OSError:
            pass

    text = extract_text(data, _filetype_from_ext(path))
    try:
        with open(cache_file, "w", encoding="utf-8") as f:
            f.write(text)
    except OSError:
        pass
    return clean_text(text)[:max_chars]


def lesson_from_upload(data_or_text: Optional[str]) -> str:
    """Accept a file path, data URL, or raw text — return cleaned text.

    Public entry point used by the API and dashboard service layer.
    """
    if not data_or_text:
        return ""
    if isinstance(data_or_text, str) and os.path.exists(data_or_text):
        try:
            return extract_from_path(data_or_text, MAX_RETURN)
        except Exception as exc:
            logger.warning("Path extraction error %s: %s", data_or_text, exc)
            return ""
    if isinstance(data_or_text, str) and data_or_text.startswith("data:"):
        try:
            header, encoded = data_or_text.split(",", 1)
            return extract_text(base64.b64decode(encoded),
                                _filetype_from_data_url_header(header))[:MAX_RETURN]
        except Exception as exc:
            logger.warning("Base64 decode failed: %s", exc)
            return ""
    try:
        return clean_text(data_or_text)[:MAX_RETURN]
    except Exception:
        return ""